from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE


SLIDE_BG = "F6F8FB"

# main dark theme -> soft light theme
FILL_COLOR_MAP = {
    "07141F": "F7F8FA",
    "0B1C2C": "F3F7FB",
    "0E2233": "EDF3F8",
    "0F2234": "EDF3F8",
    "112A40": "EAF1F7",
    "122A40": "EAF1F7",
    "2C4A6E": "D7E4F2",
    "4CB6C8": "CFEFF4",
    "7AE6F5": "E8FBFD",
    "B8703A": "F3E0D1",
    "5A8A82": "DEECE8",
    "A8864A": "F2E8D4",
    "D96A63": "F6DDD9",
}

LINE_COLOR_MAP = {
    "21405C": "D1DCE8",
    "2C4A6E": "B8C9DA",
    "DCE4EF": "DCE4EF",
}

TEXT_COLOR_MAP = {
    "FFFFFF": "1B2230",
    "A9BED4": "5A6577",
    "7AE6F5": "2C4A6E",
    "7E98B3": "7A8798",
    "B8703A": "A66A35",
    "5A8A82": "4E7B74",
    "D96A63": "BF6E68",
    "2C4A6E": "2C4A6E",
    "A8864A": "9D7C45",
    "07141F": "1B2230",
}


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def _brightness(hex_color: str) -> float:
    red = int(hex_color[0:2], 16)
    green = int(hex_color[2:4], 16)
    blue = int(hex_color[4:6], 16)
    return (0.299 * red) + (0.587 * green) + (0.114 * blue)


def _light_fill_fallback(current_hex: str) -> str:
    if _brightness(current_hex) < 55:
        return "EEF3F8"
    if _brightness(current_hex) < 100:
        return "E7EEF6"
    return current_hex


def _line_fallback(current_hex: str) -> str:
    if _brightness(current_hex) < 100:
        return "D1DCE8"
    return current_hex


def _text_fallback(current_hex: str) -> str:
    if _brightness(current_hex) > 220:
        return "1B2230"
    if _brightness(current_hex) > 160:
        return "5A6577"
    return current_hex


def recolor_fill(shape) -> None:
    try:
        fill = shape.fill
    except Exception:
        return

    if fill is None or fill.type != MSO_FILL.SOLID:
        return

    try:
        current_rgb = fill.fore_color.rgb
    except Exception:
        return

    if current_rgb is None:
        return

    current_hex = str(current_rgb)
    new_hex = FILL_COLOR_MAP.get(current_hex, _light_fill_fallback(current_hex))
    if new_hex != current_hex:
        fill.solid()
        fill.fore_color.rgb = _rgb(new_hex)


def recolor_line(shape) -> None:
    try:
        line = shape.line
    except Exception:
        return

    if line is None:
        return

    try:
        current_rgb = line.color.rgb
    except Exception:
        return

    if current_rgb is None:
        return

    current_hex = str(current_rgb)
    new_hex = LINE_COLOR_MAP.get(current_hex, _line_fallback(current_hex))
    if new_hex != current_hex:
        line.color.rgb = _rgb(new_hex)


def recolor_text(shape) -> None:
    if not getattr(shape, "has_text_frame", False):
        return

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            try:
                current_rgb = run.font.color.rgb
            except Exception:
                current_rgb = None

            if current_rgb is None:
                # keep this simple: after moving to light panels, default to dark readable text
                run.font.color.rgb = _rgb("1B2230")
                continue

            current_hex = str(current_rgb)
            new_hex = TEXT_COLOR_MAP.get(current_hex, _text_fallback(current_hex))
            if new_hex != current_hex:
                run.font.color.rgb = _rgb(new_hex)


def recolor_shape(shape) -> None:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for subshape in shape.shapes:
            recolor_shape(subshape)
        return

    recolor_fill(shape)
    recolor_line(shape)
    recolor_text(shape)


def recolor_presentation(source_path: Path, output_path: Path) -> None:
    prs = Presentation(str(source_path))

    for slide in prs.slides:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(SLIDE_BG)

        for shape in slide.shapes:
            recolor_shape(shape)

    prs.save(str(output_path))


def main() -> int:
    if len(sys.argv) < 2:
        print("Usage: python scripts/recolor_ppt_light.py <pptx_path> [output_pptx_path]")
        return 1

    ppt_path = Path(sys.argv[1])
    if not ppt_path.exists():
        print(f"File not found: {ppt_path}")
        return 1

    output_path = Path(sys.argv[2]) if len(sys.argv) >= 3 else ppt_path
    recolor_presentation(ppt_path, output_path)
    print(f"Updated light theme: {output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
